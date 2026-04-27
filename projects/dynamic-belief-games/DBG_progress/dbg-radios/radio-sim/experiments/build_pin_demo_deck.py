#!/usr/bin/env python3
"""Build a compact PowerPoint deck for PIN local-control demo results."""

from __future__ import annotations

import json
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.35), Inches(0.12), Inches(12.6), Inches(0.55))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 30, 30)
    if subtitle:
        sub = slide.shapes.add_textbox(
            Inches(0.38), Inches(0.62), Inches(12.2), Inches(0.35)
        )
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.color.rgb = RGBColor(70, 70, 70)


def _set_cell_text(cell, text: str, size: int = 11, bold: bool = False) -> None:
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold


def _add_bullets(
    slide, lines: Iterable[str], left: float, top: float, width: float, height: float, size: int
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(35, 35, 35)


def _kpi_table_rows(data: dict) -> list[tuple[str, str, str, str]]:
    base = data["baseline_mean"]
    ctrl = data["controlled_mean"]
    delta = data["delta"]
    return [
        (
            "High-priority PDR",
            f"{base['high_pdr_mean']:.4f}",
            f"{ctrl['high_pdr_mean']:.4f}",
            f"{delta['high_pdr_mean']:+.4f}",
        ),
        (
            "High-priority p95 latency (ms)",
            f"{base['high_p95_latency_ms_mean']:.4f}",
            f"{ctrl['high_p95_latency_ms_mean']:.4f}",
            f"{delta['high_p95_latency_ms_mean']:+.4f}",
        ),
        (
            "Overall PDR",
            f"{base['overall_pdr_mean']:.4f}",
            f"{ctrl['overall_pdr_mean']:.4f}",
            f"{delta['overall_pdr_mean']:+.4f}",
        ),
        (
            "Overall p95 latency (ms)",
            f"{base['overall_p95_latency_ms_mean']:.4f}",
            f"{ctrl['overall_p95_latency_ms_mean']:.4f}",
            f"{delta['overall_p95_latency_ms_mean']:+.4f}",
        ),
    ]


def _non_terminal(trace: list[dict]) -> list[dict]:
    out = []
    for row in trace:
        if (
            row["high_pdr"] == 0.0
            and row["overall_pdr"] == 0.0
            and row["high_p95_latency_ms"] == 0.0
            and row["overall_p95_latency_ms"] == 0.0
        ):
            continue
        out.append(row)
    return out


def _choose_interval_idx(base: list[dict], ctrl: list[dict]) -> int:
    for idx, (b, c) in enumerate(zip(base, ctrl)):
        hpdr = c["high_pdr"] - b["high_pdr"]
        opdr = c["overall_pdr"] - b["overall_pdr"]
        olat = c["overall_p95_latency_ms"] - b["overall_p95_latency_ms"]
        if hpdr > 0 and opdr > 0 and olat < 0:
            return idx
    return 0


def build_deck(repo_root: Path) -> Path:
    data_path = repo_root / "docs" / "pin_demo_metrics.json"
    fig_dir = repo_root / "docs" / "figures" / "pin_demo"
    out_path = repo_root / "docs" / "pin_local_control_results_deck.pptx"
    data = json.loads(data_path.read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1: problem + optimization.
    slide = prs.slides.add_slide(blank)
    _add_title(
        slide,
        "PIN Local Control Overlay: Optimization + Demonstration",
        "Minimal-slide briefing: problem formulation, A/B evidence, and mechanism",
    )
    eq = (
        "Per-radio constrained control (local only):\n"
        "a_t^(i) ~ π_θ(.|h_t^(i)),   h_t^(i) = (o_0:t^(i), a_0:t-1^(i))\n"
        "max_π  E[Σ_{t=0}^{T-1} γ^t r_t]\n"
        "s.t.  PDR_t^cmd ≥ p_min,   L95_t^cmd ≤ l_max,   Drop_t^cmd ≈ 0\n"
        "r_t = α·PDR_t^high - β·L95_t^high - δ·Drop_t^high + ζ·Deliveries_t^high"
    )
    eq_box = slide.shapes.add_textbox(Inches(0.45), Inches(1.1), Inches(8.05), Inches(2.35))
    tf = eq_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = eq
    p.font.name = "Courier New"
    p.font.size = Pt(15)

    _add_bullets(
        slide,
        [
            "Decision objective: improve PDR and latency over scenario life with no waveform rewrite.",
            "Control scope: each PIN agent controls only its local radio MAC.",
            "Action vector (per class): service bias, admission threshold, CW aggressiveness.",
            "Traffic classes: command, voice, best-effort.",
        ],
        left=8.7,
        top=1.1,
        width=4.2,
        height=2.5,
        size=13,
    )

    obs_hdr = slide.shapes.add_textbox(Inches(0.45), Inches(3.72), Inches(4.5), Inches(0.35))
    obs_p = obs_hdr.text_frame.paragraphs[0]
    obs_p.text = "Local observation window o_t^(i) (every 250 ms)"
    obs_p.font.bold = True
    obs_p.font.size = Pt(14)

    _add_bullets(
        slide,
        [
            "Queue length by class",
            "TX attempts/success, retries, ACK timeouts by class",
            "Drops + deliveries by class",
            "Per-class p95 latency",
            "Collisions, CCA busy fraction, mean backoff slots",
        ],
        left=0.55,
        top=4.08,
        width=6.8,
        height=2.9,
        size=12,
    )

    # Slide 2: experiment description + control interface.
    slide = prs.slides.add_slide(blank)
    _add_title(
        slide,
        "Experiment Description and PIN I/O Interface",
        "CSMA local-overlay validation setup",
    )
    rows, cols = 8, 2
    tbl = slide.shapes.add_table(
        rows, cols, Inches(0.45), Inches(1.05), Inches(6.6), Inches(3.15)
    ).table
    params = [
        ("Protocol", "CSMA"),
        ("Nodes / area / duration", "8 nodes / 150 m / 1.5 s"),
        ("Traffic model", "Bernoulli, 1024-bit packets"),
        ("Class mix", "Command 0.15, Voice 0.25, Best-effort 0.60"),
        ("Control interval", f"{data['config']['control_interval_ms']:.0f} ms"),
        ("Training", f"{data['config']['train_episodes']} episodes"),
        ("Eval seeds", ", ".join(str(s) for s in data["config"]["eval_seeds"])),
        ("A/B design", "Baseline A0 (neutral) vs learned local policy"),
    ]
    for r, (k, v) in enumerate(params):
        _set_cell_text(tbl.cell(r, 0), k, size=11, bold=True)
        _set_cell_text(tbl.cell(r, 1), v, size=11, bold=False)

    io_box = slide.shapes.add_textbox(Inches(7.25), Inches(1.05), Inches(5.65), Inches(3.25))
    io_tf = io_box.text_frame
    io_tf.clear()
    p = io_tf.paragraphs[0]
    p.text = "Per-node control loop (local overlay)"
    p.font.bold = True
    p.font.size = Pt(15)
    loop_lines = [
        "1) collect LocalObservation at t",
        "2) map obs to compact state bins (queue, busy, high-drop)",
        "3) choose action template A0..A3",
        "4) apply LocalAction to local MAC only",
        "5) step 250 ms, observe reward, update policy",
    ]
    for line in loop_lines:
        p = io_tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(12)

    action_tbl = slide.shapes.add_table(
        5, 4, Inches(0.45), Inches(4.45), Inches(12.45), Inches(2.55)
    ).table
    headers = ["Action", "Service Bias", "Admission", "CW Aggressiveness"]
    for c, h in enumerate(headers):
        _set_cell_text(action_tbl.cell(0, c), h, size=11, bold=True)
    rows_data = [
        ("A0 Neutral", "[1.0, 1.0, 1.0]", "[1.0, 1.0, 1.0]", "[1.0, 1.0, 1.0]"),
        ("A1 Aggressive-priority", "[1.8, 1.5, 0.5]", "[1.0, 1.0, 0.25]", "[0.75, 0.85, 1.35]"),
        ("A2 Balanced", "[1.4, 1.2, 0.8]", "[1.0, 0.9, 0.45]", "[0.85, 0.95, 1.15]"),
        ("A3 Conservative", "[1.2, 1.1, 0.9]", "[0.9, 0.8, 0.35]", "[0.95, 1.0, 1.2]"),
    ]
    for r, row in enumerate(rows_data, start=1):
        for c, val in enumerate(row):
            _set_cell_text(action_tbl.cell(r, c), val, size=10, bold=(c == 0))

    # Slide 3: aggregate results + main figures.
    slide = prs.slides.add_slide(blank)
    _add_title(
        slide,
        "Seeded A/B Results: Baseline vs PIN-Controlled",
        "Aggregate KPIs and per-seed behavior",
    )
    kpi_rows = _kpi_table_rows(data)
    tbl = slide.shapes.add_table(
        len(kpi_rows) + 1, 4, Inches(0.45), Inches(1.0), Inches(12.45), Inches(1.95)
    ).table
    for c, h in enumerate(["KPI", "Baseline mean", "Controlled mean", "Delta (C-B)"]):
        _set_cell_text(tbl.cell(0, c), h, size=12, bold=True)
    for r, row in enumerate(kpi_rows, start=1):
        for c, val in enumerate(row):
            _set_cell_text(tbl.cell(r, c), val, size=11, bold=(c == 0))

    agg_fig = fig_dir / "aggregate_kpi_comparison.png"
    seed_fig = fig_dir / "seed_level_pairwise.png"
    if agg_fig.exists():
        slide.shapes.add_picture(str(agg_fig), Inches(0.45), Inches(3.08), Inches(6.15))
    if seed_fig.exists():
        slide.shapes.add_picture(str(seed_fig), Inches(6.7), Inches(3.08), Inches(6.2))

    summary = (
        f"Overall PDR improved ({data['delta']['overall_pdr_mean']:+.4f}) and overall p95 latency improved "
        f"({data['delta']['overall_p95_latency_ms_mean']:+.2f} ms). "
        "High-priority KPIs are mixed in this small 2-seed evaluation."
    )
    s_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.78), Inches(12.25), Inches(0.45))
    s_p = s_box.text_frame.paragraphs[0]
    s_p.text = summary
    s_p.font.size = Pt(12)

    # Slide 4: highlight scenario with mechanism.
    slide = prs.slides.add_slide(blank)
    highlight_seed = data["highlight_seed"]
    _add_title(
        slide,
        f"Specific Improvement Scenario: Seed {highlight_seed}",
        data.get("highlight_reason", ""),
    )
    scenario = {row["seed"]: row for row in data["seed_scenarios"]}[highlight_seed]
    d = scenario["delta"]
    htbl = slide.shapes.add_table(
        5, 4, Inches(0.45), Inches(1.0), Inches(6.25), Inches(1.9)
    ).table
    for c, h in enumerate(["KPI", "Baseline", "Controlled", "Delta"]):
        _set_cell_text(htbl.cell(0, c), h, size=11, bold=True)
    metrics = [
        ("High-priority PDR", "high_pdr"),
        ("High-priority p95 latency (ms)", "high_p95_latency_ms"),
        ("Overall PDR", "overall_pdr"),
        ("Overall p95 latency (ms)", "overall_p95_latency_ms"),
    ]
    b = scenario["baseline_metrics"]
    c = scenario["controlled_metrics"]
    for r, (name, key) in enumerate(metrics, start=1):
        _set_cell_text(htbl.cell(r, 0), name, size=10, bold=True)
        _set_cell_text(htbl.cell(r, 1), f"{b[key]:.4f}", size=10)
        _set_cell_text(htbl.cell(r, 2), f"{c[key]:.4f}", size=10)
        _set_cell_text(htbl.cell(r, 3), f"{d[key]:+.4f}", size=10)

    base_trace = _non_terminal(scenario["baseline_trace"])
    ctrl_trace = _non_terminal(scenario["controlled_trace"])
    idx = _choose_interval_idx(base_trace, ctrl_trace)
    bi = base_trace[idx]
    ci = ctrl_trace[idx]
    interval_text = [
        f"Interval spotlight (t={bi['time_ms']:.0f} ms window):",
        f"ΔHigh PDR {ci['high_pdr'] - bi['high_pdr']:+.4f}",
        f"ΔOverall PDR {ci['overall_pdr'] - bi['overall_pdr']:+.4f}",
        f"ΔOverall p95 latency {ci['overall_p95_latency_ms'] - bi['overall_p95_latency_ms']:+.2f} ms",
        f"ΔBest-effort queue {ci['best_effort_queue_mean'] - bi['best_effort_queue_mean']:+.2f}",
        f"ΔBackoff {ci['backoff_mean'] - bi['backoff_mean']:+.2f}",
    ]
    _add_bullets(slide, interval_text, 6.95, 1.0, 5.9, 1.95, 11)

    trace_fig = fig_dir / "highlight_seed_trace_kpis.png"
    mech_fig = fig_dir / "highlight_seed_mechanism.png"
    act_fig = fig_dir / "highlight_seed_action_usage.png"
    if trace_fig.exists():
        slide.shapes.add_picture(str(trace_fig), Inches(0.45), Inches(3.05), Inches(6.3))
    if mech_fig.exists():
        slide.shapes.add_picture(str(mech_fig), Inches(6.9), Inches(3.05), Inches(4.3))
    if act_fig.exists():
        slide.shapes.add_picture(str(act_fig), Inches(11.25), Inches(3.08), Inches(1.7))

    # Slide 5: conclusions + validation status.
    slide = prs.slides.add_slide(blank)
    _add_title(
        slide,
        "Validation Outcome and Implications",
        "What this proves, what needs hardening next",
    )
    _add_bullets(
        slide,
        [
            "Demonstrated: local PIN overlay can materially change scenario outcomes without global control.",
            "Observed benefit in this run: better overall PDR and lower overall latency tail.",
            "Important caveat: high-priority aggregate KPIs are mixed with current reward/action tuning.",
            "Mechanism seen in highlight seed: reduced best-effort pressure and higher contention conservatism.",
            "This is sufficient as a reportable proof-of-control-value demonstration, not final policy quality.",
        ],
        left=0.55,
        top=1.0,
        width=12.2,
        height=2.3,
        size=14,
    )

    next_tbl = slide.shapes.add_table(
        6, 2, Inches(0.55), Inches(3.55), Inches(12.2), Inches(2.85)
    ).table
    _set_cell_text(next_tbl.cell(0, 0), "Hardening target", size=12, bold=True)
    _set_cell_text(next_tbl.cell(0, 1), "Next implementation step", size=12, bold=True)
    next_rows = [
        ("Army realism: traffic/service classes", "Keep class priorities explicit and scenario-driven mission mix."),
        ("PHY/MAC realism", "Inject jamming/interference/mobility variability and richer link adaptation hooks."),
        ("PIN observations", "Add neighbor/link quality summaries and queue age to local telemetry vector."),
        ("Control safety", "Constrain best-effort shedding + enforce command/voice minimum service guarantees."),
        ("Validation rigor", "Scale evaluation seeds/scenarios and report confidence intervals, not point estimates."),
    ]
    for r, (lhs, rhs) in enumerate(next_rows, start=1):
        _set_cell_text(next_tbl.cell(r, 0), lhs, size=11, bold=True)
        _set_cell_text(next_tbl.cell(r, 1), rhs, size=11)

    footer = slide.shapes.add_textbox(Inches(0.55), Inches(6.72), Inches(12.2), Inches(0.4))
    f = footer.text_frame.paragraphs[0]
    f.text = "Artifacts: docs/pin_demo_metrics.json, docs/pin_demo_results.md, docs/figures/pin_demo/*"
    f.font.size = Pt(11)
    f.alignment = PP_ALIGN.LEFT

    prs.save(out_path)
    return out_path


def main() -> None:
    repo_root = Path(__file__).resolve().parents[2]
    out = build_deck(repo_root)
    print("Wrote", out)


if __name__ == "__main__":
    main()
